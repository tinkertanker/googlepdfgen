import os
import sys
import time
import shlex
import shutil
import argparse
import traceback
import subprocess
import pptx
import gspread
import yaml
from pydrive2.auth import GoogleAuth
from pydrive2.drive import GoogleDrive
from alive_progress import alive_bar, alive_it


def replace_powerpoint_text(
    presentation: pptx.Presentation, replacements: dict
) -> None:
    slides = [slide for slide in presentation.slides]
    shapes = []
    for slide in slides:
        for shape in slide.shapes:
            shapes.append(shape)
    for shape in shapes:
        for match, replacement in replacements:
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            cur_text = run.text
                            new_text = cur_text.replace(str(match), str(replacement))
                            run.text = new_text
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if match in cell.text:
                            new_text = cell.text.replace(match, replacement)
                            cell.text = new_text


def main() -> None:
    try:
        # Parse arguments
        parser = argparse.ArgumentParser(
            description=
                "Generate PDFs from Google Slides/Powerpoint and Google Sheets.\n"
                "\n"
                "Refer to https://github.com/tinkertanker/googlepdfgen for more "
                "information.",
            usage=
                "python gen.py \\\n"
                "    [--config config.yaml] \\\n"
                "    [--sheet https://docs.google.com/spreadsheets/d/xxxxx/edit#gid=12345] \\\n"
                "    [--template https://docs.google.com/presentation/d/xxxxx/edit] \\\n"
                "    [--output https://drive.google.com/u/0/folders/xxxxx] \\\n"
                "    [--ppi 300]\\\n"
                "    [--soffice /Applications/LibreOffice.app/Contents/MacOS/soffice] \\\n"
                "    [--gs /opt/homebrew/bin/gs]",
            formatter_class=argparse.RawDescriptionHelpFormatter,
        )
        manual_arguments = parser.add_argument_group('manual arguments')
        manual_arguments.add_argument(
            "--config",
            help="path to YAML configuration file",
            default="config.yaml",
        )
        manual_arguments.add_argument(
            "--sheet",
            help='link to a spreadsheet on Google Sheets in the form of '
                 '"https://docs.google.com/spreadsheets/d/xxxxx/edit#gid=0" (note the '
                 'gid parameter)',
            required=False,
        )
        manual_arguments.add_argument(
            "--template",
            help='link to a Google Slides template in the form of '
                 '"https://docs.google.com/presentation/d/xxxxx/edit" OR path to a '
                 'Powerpoint template file',
            required=False,
        )
        manual_arguments.add_argument(
            "--output",
            help='link to the output folder on Google Drive in the form of '
                 '"https://drive.google.com/drive/u/0/folders/xxxxx"',
            required=False,
        )
        parser.add_argument(
            "--ppi",
            help='number of pixels per inch images should be downscaled to',
            default=300,
        )
        parser.add_argument(
            "--libreoffice",
            help="path to LibreOffice binary; defaults to $(which soffice) on "
                 "macOS/Windows and $(which libreoffice) on Linux",
            default="soffice" if sys.platform in ("win32", "darwin") else "libreoffice",
        )
        parser.add_argument(
            "--gs",
            help="path to ghostscript binary; defaults to $(which gswin32c) on Windows "
                 "and $(which gs) on macOS/Linux",
            default="gswin32c" if sys.platform in ("win32",) else "gs",
        )
        args = parser.parse_args()

        # Load configuration from YAML file if it exists
        config = {}
        if os.path.exists(args.config):
            with open(args.config, 'r') as config_file:
                config = yaml.safe_load(config_file)

        # Use command-line arguments to override YAML settings
        DATA_SHEET = args.sheet or config.get('sheet')
        TEMPLATE_SLIDES = args.template or config.get('template')
        OUTPUT_FOLDER = args.output or config.get('output')
        IMAGE_PPI = args.ppi or config.get('ppi', 300)
        LIBREOFFICE_BIN = args.libreoffice or config.get('libreoffice', "soffice" if sys.platform in ("win32", "darwin") else "libreoffice")
        GHOSTSCRIPT_BIN = args.gs or config.get('gs', "gswin32c" if sys.platform in ("win32",) else "gs")

        # Check if using config.yaml
        if args.sheet is None and args.template is None and args.output is None:
            print("Using config.yaml")

        # Validate required parameters
        if not all([DATA_SHEET, TEMPLATE_SLIDES, OUTPUT_FOLDER]):
            parser.error("Missing required parameters. Please provide them either in the YAML config file or as command-line arguments.")

        # Set up folders
        shutil.rmtree("results", ignore_errors=True)
        shutil.rmtree("_output", ignore_errors=True)
        os.makedirs("results", exist_ok=True)
        os.makedirs("_output", exist_ok=True)

        # Authenticate
        with alive_bar(
            1, title="Authenticating", bar=False, monitor=False, stats=False
        ) as bar:
            gauth = GoogleAuth(settings_file="auth.yaml")
            try:
                gauth.LocalWebserverAuth()
            except Exception as e:
                print(f"Error: {e}")
                print("Please refer to the README for instructions on how to set up your Google Cloud project and obtain the necessary credentials.")
                return
            drive = GoogleDrive(gauth)
            gc = gspread.oauth(
                credentials_filename="client_secrets.json",
                authorized_user_filename="credentials.json",
            )
            bar()

        # Fetch spreadsheet
        with alive_bar(
            1, title="Fetching spreadsheet", bar=False, monitor=False, stats=False
        ) as bar:
            data_sheet = gc.open_by_url(DATA_SHEET).get_worksheet_by_id(
                int(DATA_SHEET.split("#gid=")[-1])
            )
            records = data_sheet.get_all_records()
            records = [
                {
                    k: v
                    for k, v in record.items()
                    if k == "filename" or (k.startswith("<") and k.endswith(">"))
                }
                for record in records
            ]
            bar()

        # Fetch template slides
        with alive_bar(
            1, title="Fetching template", bar=False, monitor=False, stats=False
        ) as bar:
            if TEMPLATE_SLIDES.startswith("https://docs.google.com/presentation"):
                # TEMPLATE_SLIDE is a google slide link
                template_file = drive.CreateFile({"id": TEMPLATE_SLIDES.split("/")[-2]})
                template_file.GetContentFile(
                    "_output/template.pptx",
                    mimetype="application/vnd.openxmlformats-officedocument."
                             "presentationml.presentation",
                )
            else:
                # TEMPLATE_SLIDE is a path to a powerpoint file
                shutil.copyfile(TEMPLATE_SLIDES, "_output/template.pptx")
            bar()

        # Replace text
        for record in alive_it(records, title="Replacing text"):
            template = pptx.Presentation("_output/template.pptx")
            replace_powerpoint_text(template, record.items())
            template.save(f"_output/{record['filename']}.pptx")
        try:
            os.remove("_output/template.pptx")
        except OSError:
            pass

        # Convert to pdf
        print("Converting files to PDF (this may take a while)")
        # Workaround at https://stackoverflow.com/a/55710221
        os.system(
            f"{LIBREOFFICE_BIN} --headless --convert-to pdf --outdir _output "
            "_output/*.pptx"
        )

        # Apply compression, PDF/A compliance, and linearisation to PDFs
        for record in alive_it(records, title="Cleaning up PDFs"):

            def clean_pdf():
                command = (
                    f"{GHOSTSCRIPT_BIN} -sDEVICE=pdfwrite "
                    # Use PDF/A-2b compliance
                    "-dPDFA=2 -dPDFACompatibilityPolicy=1 "
                    # Optimise
                    "-sPDFSettings='/printer' "
                    "-sColorConversionStrategy=UseDeviceIndependentColor "
                    "-sProcessColorModel=DeviceCMYK -dEmbedAllFonts=true "
                    f"-dFastWebView=true -r{IMAGE_PPI} "
                    # Set input and output
                    f'-q -o "results/{record["filename"]}.pdf" '
                    f'"_output/{record["filename"]}.pdf"'
                )
                try:
                    subprocess.check_output(
                        shlex.split(command), stderr=subprocess.STDOUT
                    )
                except subprocess.CalledProcessError as e:
                    # Sometimes weird things happen, so we'll just try again
                    print(
                        f"Retrying `{command}` due to error code {e.returncode}",
                    )
                    clean_pdf()

            clean_pdf()

        # Upload to Google Drive
        file_links = []
        for record in alive_it(records, title="Uploading PDFs to Google Drive"):

            def upload_file():
                try:
                    file = drive.CreateFile(
                        {
                            "title": f"{record['filename']}.pdf",
                            "parents": [
                                {
                                    "kind": "drive#parentReference",
                                    "id": OUTPUT_FOLDER.split("/")[-1],
                                }
                            ],
                        }
                    )
                    file.SetContentFile(f"results/{record['filename']}.pdf")
                    file.Upload()
                    file_links.append(file["alternateLink"])
                except KeyboardInterrupt:
                    raise
                except:
                    time.sleep(1)
                    # All sorts of network errors may occur, so we'll just try again
                    print(
                        f"Retrying upload of results/{record['filename']}.pdf "
                        f"due to the following exception:\n{traceback.format_exc()}",
                    )
                    upload_file()

            upload_file()

        # Update file links in spreadsheet
        with alive_bar(
            1,
            title="Updating file links in spreadsheet",
            bar=False,
            monitor=False,
            stats=False,
        ) as bar:

            def update_file_links():
                try:
                    file_col_index = [
                        cell
                        for cell in data_sheet.get_all_cells()
                        if cell.row == 1 and cell.value == "file"
                    ][0].col
                    file_col_id = ""
                    while file_col_index > 0:
                        file_col_index, remainder = divmod(file_col_index - 1, 26)
                        file_col_id = chr(ord("A") + remainder) + file_col_id
                    data_sheet.update(
                        f"{file_col_id}2", [file_links], major_dimension="COLUMNS"
                    )
                except KeyboardInterrupt:
                    raise
                except:
                    time.sleep(1)
                    # All sorts of network errors may occur, so we'll just try again
                    print(
                        f"Retrying update of {file_col_id}2 due to the "
                        f"following exception:\n{traceback.format_exc()}",
                    )
                    update_file_links()

            update_file_links()

            bar()

        # Prompt user to make the folder public
        make_public = input(f"\nDo you want to make the output folder public? (y/n): ").lower().strip() == 'y'
        
        if make_public:
            with alive_bar(1, title="Making folder public", bar=False, monitor=False, stats=False) as bar:
                try:
                    folder = drive.CreateFile({'id': OUTPUT_FOLDER.split("/")[-1]})
                    folder.InsertPermission({
                        'type': 'anyone',
                        'role': 'reader',
                        'withLink': True
                    })
                    print(f"\nThe folder has been made public. Anyone with the link can view it.")
                except Exception as e:
                    print(f"\nError making folder public: {e}")
                bar()

        print(f"\nYou can access your output folder at: {OUTPUT_FOLDER}")

    except KeyboardInterrupt:
        pass
    finally:
        # Clean up
        print()
        with alive_bar(
            1, title="Cleaning up", bar=False, monitor=False, stats=False
        ) as bar:
            shutil.rmtree("_output", ignore_errors=True)
            bar()


if __name__ == "__main__":
    main()
